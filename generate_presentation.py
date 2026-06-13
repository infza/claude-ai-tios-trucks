#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Claude AI: Tipy a Triky - PowerPoint Prezentácia Generator
Tento skript generuje kompletný PowerPoint súbor s prezentáciou
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Pridaj titulnú snímku"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Prázdny layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 31, 31)  # Temnú farbu
    
    # Titulka
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Podtitul
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, is_two_column=False):
    """Pridaj obsah snímku"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 31, 31)
    
    # Titulka
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(66, 175, 250)
    
    # Obsah
    if is_two_column:
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(5))
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5))
        
        for box, column_content in zip([left_box, right_box], [content_items[0], content_items[1]]):
            text_frame = box.text_frame
            text_frame.word_wrap = True
            for i, item in enumerate(column_content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.space_before = Pt(10)
                p.space_after = Pt(10)
                p.level = 0
    else:
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            p.level = 0
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Pridaj snímku s tabuľkou"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 31, 31)
    
    # Titulka
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(66, 175, 250)
    
    # Tabuľka
    rows_count = len(rows) + 1
    cols_count = len(headers)
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(8.4)
    height = Inches(4.5)
    
    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
    
    # Header
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(66, 175, 250)
        
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.size = Pt(18)
    
    # Riadky
    for row_idx, row_data in enumerate(rows, 1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(50, 50, 50)
            else:
                cell.fill.fore_color.rgb = RGBColor(40, 40, 40)
            
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.size = Pt(16)
    
    return slide

def create_presentation():
    """Vytvor kompletný prezentáciu"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Titulka
    add_title_slide(prs, "🤖 Claude AI", "Tipy a Triky")
    
    # Slide 2: Obsah
    add_content_slide(prs, "📑 Obsah Prezentácie", [
        "✨ Úvod do Claude AI",
        "💡 Základné Tipy",
        "🎯 Pokročilé Triky",
        "🎨 Promptovanie",
        "📊 Prípadové Štúdie",
        "✅ Best Practices",
        "🔒 Bezpečnosť",
        "🚀 Pokročilé Techniky"
    ])
    
    # Slide 3: Úvod - Čo je Claude AI?
    add_content_slide(prs, "🤖 Čo je Claude AI?", [
        "✓ Pokročilý jazykový model vyvinutý Anthropic",
        "✓ Vysoká kvalita konverzácií",
        "✓ Zabudované etické princípy",
        "✓ Všestrannosť - od kódovania po tvorivé písanie",
        "✓ Dlhý kontext - schopnosť pracovať s dlhými textami"
    ])
    
    # Slide 4: Verzie Claude
    add_table_slide(prs, "📦 Verzie Claude AI", 
        ["Verzia", "Charakteristiky", "Ideálne na"],
        [
            ["Opus", "Najpokročilejší", "Zložité úlohy, analýza"],
            ["Sonnet", "Rýchly a efektívny", "Bežné úlohy"],
            ["Haiku", "Ľahký a rýchly", "Jednoduché operácie"]
        ])
    
    # Slide 5: Základný Tip 1
    add_content_slide(prs, "💡 Tip 1: Jasné Otázky", [
        "❌ Zlé: 'Pomôž mi s kódovaním.'",
        "",
        "✅ Dobré: 'Pomôž mi napísať Python funkciu, ktorá čita CSV súbor a vracia priemernú hodnotu numerického stĺpca price.'"
    ])
    
    # Slide 6: Základný Tip 2
    add_content_slide(prs, "💡 Tip 2: Kontextové Informácie", [
        "🔧 Technologický stack",
        "📦 Verzia nástrojov",
        "🎯 Cieľ úlohy",
        "📤 Očakávaný výstup",
        "⚙️ Obmedzenia a požiadavky"
    ])
    
    # Slide 7: Základný Tip 3
    add_content_slide(prs, "💡 Tip 3: Formátovanie Výstupu", [
        "📋 'Poskytnúť ako JSON'",
        "📊 'Formátovať ako tabuľka'",
        "📝 'Výstup ako Markdown'",
        "🔢 'Číselný zoznam'",
        "📑 'Polia s nadpismi'"
    ])
    
    # Slide 8: Základný Tip 4
    add_content_slide(prs, "💡 Tip 4: Postupné Spresnenie", [
        "✓ 'Áno, ale pridaj viac príkladov'",
        "✓ 'Zmeň to na Go namiesto Python'",
        "✓ 'Zameraj sa viac na bezpečnosť'",
        "✓ 'Urob to jednoduchšie'",
        "✓ 'Pridaj komentáre do kódu'"
    ])
    
    # Slide 9: Základný Tip 5
    add_content_slide(prs, "💡 Tip 5: Rozdeľ Zložité Úlohy", [
        "❌ Zle: Vytvor kompletný e-shop backend",
        "",
        "✅ Dobre:",
        "  1. API endpoint autentifikácie",
        "  2. Endpointy pre produkty",
        "  3. Integrácia platieb"
    ])
    
    # Slide 10: Pokročilý Trik 1
    add_content_slide(prs, "🎯 Trik 1: Role-Based Prompting", [
        "Definuj rolu pre Claude:",
        "",
        "'Si skúsený Python vývojár s 10 rokmi skúsenosti.",
        "Pomôž mi optimalizovať tento kód na výkon.'"
    ])
    
    # Slide 11: Pokročilý Trik 2
    add_content_slide(prs, "🎯 Trik 2: Extended Thinking", [
        "Pre zložité problémy:",
        "",
        "'Prosím, premysli si tento problém krok za krokom",
        "a potom mi poskytni kompletné riešenie.'"
    ])
    
    # Slide 12: Pokročilý Trik 3
    add_content_slide(prs, "🎯 Trik 3: Iteratívne Zlepšovanie", [
        "📝 Počiatočný draft",
        "🔍 'Zameraj sa viac na X'",
        "📚 'Pridaj príklady'",
        "✂️ 'Zjednoduš to'",
        "🎨 'Vylepši formátovanie'"
    ])
    
    # Slide 13: Pokročilý Trik 4
    add_content_slide(prs, "🎯 Trik 4: Multi-Turn Konverzácie", [
        "Udržiavaj kontext v rámci rozhovoru:",
        "",
        "1. Objasnite princípy X",
        "2. Ako sa to vzťahuje na Y?",
        "3. Môžeš mi to ukázať na príklade?"
    ])
    
    # Slide 14: Pokročilý Trik 5-7
    add_content_slide(prs, "🎯 Pokročilé Triky", [
        "5️⃣ Prompt Chaining - spájaj prompty v sérii",
        "",
        "6️⃣ Few-Shot Learning - poskytnú príklady",
        "",
        "7️⃣ Chain-of-Thought - vysvetluj myslenie"
    ])
    
    # Slide 15: Štruktúra Prompta
    add_content_slide(prs, "🎨 Štruktúra Efektívneho Prompta", [
        "[ROLA] - Kto si?",
        "[KONTEXT] - Aká je situácia?",
        "[ÚLOHA] - Čo máš robiť?",
        "[FORMÁT] - Ako to chceš?",
        "[PRÍKLADY] - Prípadne príklady",
        "[OBMEDZENIA] - Aké sú limity?"
    ])
    
    # Slide 16: Bezpečné Slovníčky
    add_table_slide(prs, "🎨 Bezpečné Slovníčky",
        ["Termín", "Efekt"],
        [
            ["prosím", "Zlepšuje tón"],
            ["krok za krokom", "Podrobnejšie vysvetlenia"],
            ["ako začiatočník", "Jednoduchšie vysvetlenia"],
            ["porovnaj", "Viacero možností"],
            ["prečo", "Hlbšie vysvetlenie"]
        ])
    
    # Slide 17: Prípadová štúdia 1
    add_content_slide(prs, "📊 Prípadová Štúdia 1: Aplikácia", [
        "Problém: Potrebujem architektúru React aplikácie",
        "",
        "Riešenie: Navrhni architektúru s požiadavkami",
        "- 50+ produktov",
        "- Nákupný košík",
        "- Viacerí používatelia",
        "- SEO optimalizácia"
    ])
    
    # Slide 18: Prípadová štúdia 2
    add_content_slide(prs, "📊 Prípadová Štúdia 2: Ladenie Kódu", [
        "Problém: Čudný bug v produkcii",
        "",
        "Riešenie: Analyzuj stack trace s kontextom",
        "- Node.js 18",
        "- Express app",
        "- PostgreSQL databáza",
        "- Problém sa vyskytuje náhodne"
    ])
    
    # Slide 19: Prípadová štúdia 3
    add_content_slide(prs, "📊 Prípadová Štúdia 3: Dokumentácia", [
        "Problém: Potrebujem API dokumentáciu",
        "",
        "Riešenie: Vytvor OpenAPI 3.0 špecifikáciu",
        "- GET, POST, PUT, DELETE endpointy",
        "- YAML formát",
        "- Príklady a operácie"
    ])
    
    # Slide 20: Best Practices - Bezpečnosť
    add_content_slide(prs, "✅ Best Practice 1: Bezpečnosť", [
        "✅ ROBIŤ:",
        "  • Nedeliť citlivé informácie",
        "  • Anonymizovať dáta",
        "  • Overovať citlivé výstupy",
        "",
        "❌ NEROBIŤ:",
        "  • Nedeliť API kľúče",
        "  • Nedeliť hesla",
        "  • Nedeliť obchodné tajomstvá"
    ])
    
    # Slide 21: Best Practices - Efektivita
    add_table_slide(prs, "✅ Best Practice 2: Efektivita",
        ["Typ Úlohy", "Odporúčaný Model"],
        [
            ["Rýchle otázky", "Sonnet alebo Haiku"],
            ["Zložité úlohy", "Opus"],
            ["Dlhý kontext", "Sonnet alebo Opus"],
            ["Nákladosť", "Zvážiť model vs. zložitosť"]
        ])
    
    # Slide 22: Best Practices - Kvalita
    add_content_slide(prs, "✅ Best Practice 3: Kvalita Výstupu", [
        "1️⃣ Jasne definuj očakávania",
        "2️⃣ Poskytnú kontext",
        "3️⃣ Požiadaj o vysvetlenia",
        "4️⃣ Iteruj na spätnej väzbe",
        "5️⃣ Overuj výstupy"
    ])
    
    # Slide 23: Bezpečnosť
    add_content_slide(prs, "🔒 Bezpečnosť Claude AI", [
        "✅ Nedeliť hesla alebo kľúče",
        "✅ Nedeliť osobné údaje",
        "✅ Nedeliť obchodné tajomstvá",
        "✅ Veriť slepé bez overenia",
        "✅ Ignorovať bezpečnostné varovania"
    ])
    
    # Slide 24: Verifikácia
    add_content_slide(prs, "🔒 Verifikácia Výstupu", [
        "🧪 Testuj kódové príklady",
        "✓ Overi faktické tvrdenia",
        "🔍 Kontroluj bezpečnostné obavy",
        "⚠️ Validuj citlivé výstupy",
        "👥 Konzultuj s expertmi ak je potrebné"
    ])
    
    # Slide 25: Pokročilé Techniky
    add_content_slide(prs, "🚀 Pokročilé Techniky", [
        "1️⃣ Rubber Duck Debugging - vysvetli problém",
        "2️⃣ Analogické Myslenie - porovnávaj s jednoduchším",
        "3️⃣ Brainstorming - generuj nápady",
        "4️⃣ Analýza Konkurencie - porovnávaj produkty"
    ])
    
    # Slide 26: Kľúčové Poznatky
    add_content_slide(prs, "🎯 Kľúčové Poznatky", [
        "💬 Jasnosť je kľúč - presné otázky = presné odpovede",
        "📚 Kontext je dôležitý - čím viac info, tým lepšie",
        "🔄 Iterácia funguje - prvá odpoveď nie je vždy ideálna",
        "✓ Overuj výstupy - myslite kriticky",
        "🧪 Experimentuj - nájdi techniky, ktoré ti fungujú"
    ])
    
    # Slide 27: Ďalšie Kroky
    add_content_slide(prs, "📝 Ďalšie Kroky", [
        "1️⃣ Skús Claude AI na svojej príštej úlohe",
        "2️⃣ Experimentuj s rôznymi promptami",
        "3️⃣ Zaznamenávaj, čo funguje",
        "4️⃣ Zdieľaj svoje poznatky s tímom",
        "5️⃣ Neustále sa učiť a zlepšovať"
    ])
    
    # Slide 28: Záver
    add_title_slide(prs, "Ďakujem! 🙏", "Teraz tvoj rad!")
    
    # Ulož prezentáciu
    output_file = "Claude_AI_Tipy_a_Triky.pptx"
    prs.save(output_file)
    print(f"✅ Prezentácia bola úspešne vytvorená: {output_file}")
    return output_file

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Chyba: Balíček 'python-pptx' nie je nainštalovaný.")
        print("Prosím, nainštaluj ho pomocou:")
        print("  pip install python-pptx")
    except Exception as e:
        print(f"❌ Chyba: {str(e)}")
